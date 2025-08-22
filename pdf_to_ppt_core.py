# file: pdf_to_ppt_core.py (修正版，去阴影)

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import io
import os
import logging
import re
from collections import defaultdict

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def convert_pt_to_inches(pt_value):
    """将点(pt)转换为英寸"""
    return pt_value / 72.0  # 1 pt = 1/72 inch

def convert_pt_to_emu(pt_value):
    """将点(pt)转换为EMU"""
    return Emu(pt_value * 12700)  # 1 pt = 12700 EMU

def is_shadow_image(img_rect, other_rects, tolerance=2):
    """
    检测是否为阴影图片（与其他图片重叠且尺寸稍大）
    :param img_rect: 当前图片的矩形区域
    :param other_rects: 其他图片的矩形区域列表
    :param tolerance: 尺寸差异容差（pt）
    :return: 如果是阴影返回True，否则返回False
    """
    for other_rect in other_rects:
        # 检查是否重叠
        if (img_rect.x0 <= other_rect.x1 and img_rect.x1 >= other_rect.x0 and
            img_rect.y0 <= other_rect.y1 and img_rect.y1 >= other_rect.y0):
            
            # 检查当前图片是否比另一个图片稍大（可能是阴影）
            width_diff = (img_rect.width - other_rect.width)
            height_diff = (img_rect.height - other_rect.height)
            
            # 如果当前图片比另一个图片稍大，且重叠面积超过80%，则认为是阴影
            if (width_diff > tolerance and height_diff > tolerance and
                width_diff < 20 and height_diff < 20):  # 阴影通常不会太大
                
                overlap_area = (min(img_rect.x1, other_rect.x1) - max(img_rect.x0, other_rect.x0)) * \
                              (min(img_rect.y1, other_rect.y1) - max(img_rect.y0, other_rect.y0))
                img_area = img_rect.width * img_rect.height
                overlap_ratio = overlap_area / img_area
                
                if overlap_ratio > 0.8:
                    return True
    
    return False

def pdf_to_ppt(pdf_path, pptx_path):
    """
    将PDF转换为PPTX，保持布局且文本可编辑
    :param pdf_path: PDF文件路径
    :param pptx_path: 输出的PPTX文件路径
    :return: 成功返回True，失败返回False
    """
    try:
        # 打开PDF文件
        doc = fitz.open(pdf_path)
        prs = Presentation()
        
        # 获取PDF页面尺寸
        first_page = doc.load_page(0)
        pdf_width = first_page.rect.width
        pdf_height = first_page.rect.height
        
        # 设置PPT页面尺寸（保持比例）
        # 将PDF的pt单位转换为英寸
        pdf_width_inches = convert_pt_to_inches(pdf_width)
        pdf_height_inches = convert_pt_to_inches(pdf_height)
        
        # 设置PPT页面尺寸
        ppt_width = Inches(pdf_width_inches)
        ppt_height = Inches(pdf_height_inches)
        
        prs.slide_width = ppt_width
        prs.slide_height = ppt_height
        
        logger.info(f"开始转换: {pdf_path}")
        logger.info(f"PDF尺寸: {pdf_width}x{pdf_height} pt")
        logger.info(f"PPT尺寸: {ppt_width.inches:.2f}x{ppt_height.inches:.2f} 英寸")
        
        # 计算缩放比例（用于坐标转换）
        scale_factor = 1.0  # 因为PPT尺寸已经按比例设置，所以缩放比例为1
        
        for page_num in range(len(doc)):
            try:
                page = doc.load_page(page_num)
                # 创建空白幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                logger.info(f"处理第 {page_num + 1} 页...")
                
                # 1. 首先提取文本块并精确定位
                text_dict = page.get_text("dict")
                
                for block in text_dict.get("blocks", []):
                    if block["type"] == 0:  # 文本块
                        try:
                            # 获取文本块的边界框和文本内容
                            bbox = block["bbox"]
                            x0, y0, x1, y1 = bbox
                            
                            # 转换为PPT坐标（英寸）
                            left = Inches(convert_pt_to_inches(x0))
                            top = Inches(convert_pt_to_inches(y0))
                            width = Inches(convert_pt_to_inches(x1 - x0))
                            height = Inches(convert_pt_to_inches(y1 - y0))
                            
                            # 提取文本内容
                            text_content = ""
                            font_size = Pt(12)  # 默认值
                            font_name = "Arial"
                            
                            for line in block.get("lines", []):
                                for span in line.get("spans", []):
                                    text_content += span["text"]
                                    # 获取字体信息（使用第一个span的字体）
                                    if "size" in span and font_size == Pt(12):
                                        font_size = Pt(span["size"])
                                    if "font" in span and font_name == "Arial":
                                        font_name = span["font"]
                                text_content += "\n"  # 保留换行
                            
                            text_content = text_content.strip()
                            if text_content:
                                # 添加文本框
                                txBox = slide.shapes.add_textbox(left, top, width, height)
                                tf = txBox.text_frame
                                tf.text = text_content
                                tf.word_wrap = False  # 禁用自动换行
                                
                                # 设置文本格式
                                for paragraph in tf.paragraphs:
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = font_size
                                        run.font.name = font_name
                                        run.font.color.rgb = 0x000000  # 黑色文字
                                
                                logger.info(f"添加文本: '{text_content[:30]}...' 位置: {left.inches:.2f}, {top.inches:.2f}")
                                
                        except Exception as e:
                            logger.error(f"处理文本块时出错: {e}")
                            continue
                
                # 2. 处理图片（修复阴影问题）
                try:
                    img_list = page.get_images()
                    img_rects = []  # 存储所有图片的矩形区域
                    
                    # 首先收集所有图片的矩形区域
                    for img in img_list:
                        xref = img[0]
                        rects = page.get_image_rects(xref)
                        if rects:
                            img_rects.extend(rects)
                    
                    # 然后处理图片，跳过阴影图片
                    processed_images = set()  # 记录已处理的图片xref
                    
                    for img_index, img in enumerate(img_list):
                        xref = img[0]
                        if xref in processed_images:
                            continue
                            
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # 获取图片位置
                        current_rects = page.get_image_rects(xref)
                        if not current_rects:
                            continue
                            
                        img_rect = current_rects[0]
                        
                        # 检查是否为阴影（与其他图片重叠且尺寸稍大）
                        other_rects = [r for r in img_rects if r != img_rect]
                        if is_shadow_image(img_rect, other_rects):
                            logger.info(f"跳过阴影图片 {img_index + 1}")
                            continue
                        
                        # 转换为PPT坐标
                        left = Inches(convert_pt_to_inches(img_rect.x0))
                        top = Inches(convert_pt_to_inches(img_rect.y0))
                        width = Inches(convert_pt_to_inches(img_rect.width))
                        height = Inches(convert_pt_to_inches(img_rect.height))
                        
                        # 添加图片
                        slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
                        logger.info(f"添加图片 {img_index + 1}")
                        
                        processed_images.add(xref)
                            
                except Exception as e:
                    logger.error(f"处理图片时出错: {e}")
                
                # 3. 添加页码
                try:
                    footer_left = Inches(0)
                    footer_top = ppt_height - Inches(0.5)
                    footer_width = ppt_width
                    footer_height = Inches(0.5)
                    
                    footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
                    footer_frame = footer_box.text_frame
                    footer_frame.text = f"第 {page_num + 1} 页 / 共 {len(doc)} 页"
                    
                    for paragraph in footer_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            run.font.color.rgb = 0x666666
                    
                except Exception as e:
                    logger.error(f"添加页码时出错: {e}")
                    
            except Exception as e:
                logger.error(f"处理第 {page_num + 1} 页时出错: {e}")
                continue
        
        # 保存文件
        prs.save(pptx_path)
        doc.close()
        
        logger.info(f"转换完成: {pdf_path} -> {pptx_path}")
        return True
        
    except Exception as e:
        logger.error(f"转换过程中出错: {e}")
        import traceback
        traceback.print_exc()
        return False

# 调试函数
def debug_pdf_structure(pdf_path):
    """调试PDF结构"""
    try:
        doc = fitz.open(pdf_path)
        page = doc.load_page(0)
        
        print("=== PDF结构分析 ===")
        print(f"页数: {len(doc)}")
        print(f"页面尺寸: {page.rect.width} x {page.rect.height} pt")
        print(f"页面尺寸(英寸): {convert_pt_to_inches(page.rect.width):.2f} x {convert_pt_to_inches(page.rect.height):.2f} 英寸")
        
        # 获取文本块信息
        text_dict = page.get_text("dict")
        print(f"文本块数量: {len(text_dict.get('blocks', []))}")
        
        for i, block in enumerate(text_dict.get("blocks", [])):
            if block["type"] == 0:  # 文本块
                print(f"\n文本块 {i}: 位置 {block['bbox']}")
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        print(f"  文本: '{span['text']}', 字体: {span.get('font', '未知')}, 大小: {span.get('size', '未知')}pt")
        
        # 分析图片信息
        img_list = page.get_images()
        print(f"\n图片数量: {len(img_list)}")
        for i, img in enumerate(img_list):
            xref = img[0]
            rects = page.get_image_rects(xref)
            if rects:
                rect = rects[0]
                print(f"图片 {i + 1}: 位置 ({rect.x0:.1f}, {rect.y0:.1f}) - ({rect.x1:.1f}, {rect.y1:.1f}), "
                      f"尺寸 {rect.width:.1f}x{rect.height:.1f} pt")
        
        doc.close()
        
    except Exception as e:
        print(f"调试失败: {e}")

# 表格分析函数（空函数，保持接口兼容）
def analyze_table_structure(pdf_path, page_num=0):
    """分析PDF中的表格结构（兼容接口）"""
    print("表格分析功能暂未实现")
    return []
